# -*- coding: utf-8 -*-
"""ONE-SHOT patch (do not re-run): SDD Appendix B (AI/ML Scope) + client deck AI slide (inserted before Benefits)."""
from docx import Document
from docx.shared import Pt, Inches, RGBColor as DRGB
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import copy

# ============ 1. SDD Appendix B ============
SDD = r"C:\Users\jnamm\OneDrive\Desktop\Service Cloud for Utilities\Enrollment Center\Enrollment Center - SAP Service Cloud V2 Solution Design v2.0.docx"
doc = Document(SDD)
already = any(p.text.startswith("Appendix B") for p in doc.paragraphs)
if not already:
    doc.add_page_break()
    doc.add_heading("Appendix B — AI/ML Scope", level=1)
    def para(text, bold=False):
        pr = doc.add_paragraph(); r = pr.add_run(text); r.bold = bold
        r.font.name = "Arial"; r.font.size = Pt(10.5)
        pr.paragraph_format.space_after = Pt(6)
        return pr
    para("The solution is extended by an AI/ML layer governed by one principle: ML recommends — BRFplus "
         "decides. Hard eligibility, warnings and pre-active behavior remain deterministic decision-table "
         "outcomes; AI operates above that line with the agent in the loop. Full detail (use cases AI-1…AI-7, "
         "technology mapping, phasing, governance) is maintained in 'Enrollment Center – AI/ML Roadmap v1.0'.")
    HEADER_FILL="1F4E79"; ALT_FILL="DEEAF6"
    def shade(cell, fill):
        tcPr = cell._tc.get_or_add_tcPr(); shd = OxmlElement('w:shd')
        shd.set(qn('w:val'),'clear'); shd.set(qn('w:fill'),fill); tcPr.append(shd)
    rows = [
     ["AI-1 Next-best-program recommendation","Propensity ranking of Available Programs with reasons (EC-10 scoring API; HANA Cloud PAL)","Phase 1"],
     ["AI-2 Grounded agent copilot","Plain-language program explanations and eligibility remediation, grounded on the program catalog (Generative AI Hub / Joule Studio; CX AI Toolkit)","Phase 1"],
     ["AI-3 Bill-impact simulation","Recompute interval data under the target rate before enrolling — recommend only when beneficial","Phase 2"],
     ["AI-4 Payment-plan default risk","Predict DPA/AMP breakage; proactive outreach","Phase 2"],
     ["AI-5 Assistance document intelligence","Income-proof extraction (Document Information Extraction) for LIHEAP/PIPP/CARE","Phase 2"],
     ["AI-6 / AI-7","Exception triage learning; DER baseline & fleet forecasting (DER platform side)","Phase 3"],
    ]
    t = doc.add_table(rows=1, cols=3); t.style = "Table Grid"; t.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i,h in enumerate(["Use Case","Summary","Phase"]):
        c = t.rows[0].cells[i]; c.text=""
        r = c.paragraphs[0].add_run(h); r.bold=True; r.font.size=Pt(9); r.font.name="Arial"
        r.font.color.rgb = DRGB(0xFF,0xFF,0xFF); shade(c, HEADER_FILL)
    for ri,row in enumerate(rows):
        cells = t.add_row().cells
        for i,val in enumerate(row):
            cells[i].text=""
            run = cells[i].paragraphs[0].add_run(val); run.font.size=Pt(9); run.font.name="Arial"
            if ri%2==1: shade(cells[i], ALT_FILL)
    for i,w in enumerate([1.7,3.6,1.2]):
        for row in t.rows: row.cells[i].width = Inches(w)
    para("The interactive demo previews AI-1 and AI-2 (mock scores, grounded explanation composed from "
         "catalog and customer facts), clearly labeled as AI and outside the eligibility decision path.")
    doc.save(SDD)
    print("SDD Appendix B added")
else:
    print("SDD already has Appendix B - skipped")

# ============ 2. Deck: AI slide ============
DECK = r"C:\Users\jnamm\OneDrive\Desktop\Service Cloud for Utilities\Enrollment Center\Enrollment Center - Client Presentation v1.0.pptx"
NAVY=RGBColor(0x12,0x23,0x3F); BLUE=RGBColor(0x00,0x70,0xF2); BLUED=RGBColor(0x00,0x57,0xD2)
AMBER=RGBColor(0xFF,0xB3,0x00); AMBERT=RGBColor(0x3B,0x2A,0x00); WHITE=RGBColor(0xFF,0xFF,0xFF)
INK=RGBColor(0x1D,0x2D,0x3E); MUTED=RGBColor(0x55,0x6B,0x82); LIGHT=RGBColor(0xF5,0xF6,0xF7)
CARD_B=RGBColor(0xD9,0xD9,0xD9); PURP=RGBColor(0x6C,0x32,0xA9); PURPC=RGBColor(0xF0,0xEA,0xFD)
GREENC=RGBColor(0xEB,0xF5,0xEB); GREEN=RGBColor(0x18,0x89,0x18); BLUEC=RGBColor(0xE3,0xF0,0xFF)

prs = Presentation(DECK)
if any("AI on top of the rules" in (sh.text_frame.text if sh.has_text_frame else "") for sl in prs.slides for sh in sl.shapes):
    print("deck already has AI slide - skipped")
else:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    def tb(x,y,w,h,text,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT,ls=1.0):
        box=s.shapes.add_textbox(PIn(x),PIn(y),PIn(w),PIn(h)); tf=box.text_frame; tf.word_wrap=True
        tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
        for i,ln in enumerate(text.split("\n")):
            p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
            if ls!=1.0: p.line_spacing=ls
            r=p.add_run(); r.text=ln; r.font.size=PPt(size); r.font.bold=bold; r.font.name="Arial"; r.font.color.rgb=color
        return box
    def pill(x,y,w,h,text,fill,tc,size=11.5):
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PIn(x),PIn(y),PIn(w),PIn(h))
        sh.adjustments[0]=0.5; sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.fill.background(); sh.shadow.inherit=False
        tf=sh.text_frame; tf.margin_left=PIn(0.02); tf.margin_right=PIn(0.02); tf.margin_top=0; tf.margin_bottom=0
        p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=text; r.font.size=PPt(size); r.font.bold=True; r.font.name="Arial"; r.font.color.rgb=tc
    def card(x,y,w,h,fill=LIGHT,border=CARD_B):
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PIn(x),PIn(y),PIn(w),PIn(h))
        sh.adjustments[0]=0.08; sh.fill.solid(); sh.fill.fore_color.rgb=fill
        sh.line.color.rgb=border; sh.line.width=PPt(0.75); sh.shadow.inherit=False
    tb(0.6,0.38,12.1,0.75,"AI on top of the rules",size=30,bold=True)
    tb(0.6,1.02,12.1,0.4,"Machine learning recommends — the rules engine decides. Full portfolio in the AI/ML Roadmap.",size=14,color=MUTED)
    cards=[("AI-1 · Recommend","Next-best-program ranking with reasons — agents pitch the right program first (live in the demo)",PURPC,PURP),
           ("AI-2 · Explain","Plain-language program explanations for THIS customer, grounded on the catalog (live in the demo)",BLUEC,BLUED),
           ("AI-3 · Verify savings","Bill-impact simulation on interval data before rate enrollment — recommend only when it saves",GREENC,GREEN),
           ("AI-4 · Protect","Payment-plan default prediction, document intelligence for assistance programs, exception triage",AMBER,AMBERT)]
    xs=[0.6,6.75]; ys=[1.8,3.55]; idx=0
    for yy in ys:
        for xx in xs:
            h,b,f,c=cards[idx]; idx+=1
            card(xx,yy,5.95,1.6)
            pill(xx+0.3,yy+0.22,2.3,0.34,h,f,c)
            tb(xx+0.3,yy+0.68,5.4,0.85,b,size=12,color=MUTED,ls=1.12)
    card(0.6,5.4,12.1,1.35,fill=NAVY,border=NAVY)
    tb(0.9,5.6,11.5,0.95,"Governance: hard eligibility stays in auditable BRFplus decision tables — no model can enroll, reject or override. "
       "AI ranks, explains, predicts and drafts, always labeled, always with the agent in the loop.",size=13,color=RGBColor(0xCA,0xDC,0xFC),ls=1.15)
    # move slide from last position to before the Benefits slide (index 10, 0-based)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[-1])
    sldIdLst.insert(10, ids[-1])
    prs.save(DECK)
    print("deck AI slide inserted at position 11")
