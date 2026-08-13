# -*- coding: utf-8 -*-
"""ONE-SHOT patch (do not re-run): update client deck (personas count, demo URLs) and
SDD v2.0 (device/DER context note in UI section) in place.
The deck's original generator was lost to a scratchpad wipe; the deck is maintained in place from now on."""
from pptx import Presentation
from docx import Document
from docx.shared import Pt
from docx.oxml import OxmlElement
from docx.text.paragraph import Paragraph

# ---------- 1. Client deck ----------
DECK = r"C:\Users\jnamm\OneDrive\Desktop\Service Cloud for Utilities\Enrollment Center\Enrollment Center - Client Presentation v1.0.pptx"
REPL = [
 ("4 customer personas, 28 programs", "7 eligibility personas · 28 programs"),
 ("nammi2011.github.io/enrollment-center-demo/", "demo.svgbs.com  (pending DNS; live: sri-varahi-global-business-services-llc.github.io/enrollment-center-demo)"),
 ("…/enrollment-center-demo/workcenter.html", "demo.svgbs.com/workcenter.html"),
 ("…/enrollment-center-demo/movein.html", "demo.svgbs.com/movein.html"),
 ("Demo: nammi2011.github.io/enrollment-center-demo   ·   Full design set available for review",
  "Demo: demo.svgbs.com   ·   Full design set available for review"),
 ("Switch personas — credit, medical,\ncommercial rules change live",
  "Switch among 7 personas — non-AMI,\nmedical, prepay, commercial rules change live"),
]
prs = Presentation(DECK)
count = 0
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new in REPL:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        count += 1
                    # multi-line runs were written as separate runs; also try line-by-line
                    elif "\n" not in old and old in run.text:
                        pass
# handle the two-line persona string written as separate runs
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame: continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip() == "Switch personas — credit, medical,":
                    run.text = "Switch among 7 personas — non-AMI,"; count += 1
                if run.text.strip() == "commercial rules change live":
                    run.text = "medical, prepay, commercial — rules change live"; count += 1
prs.save(DECK)
print("deck replacements:", count)

# ---------- 2. SDD v2.0: device/DER note in UI section ----------
SDD = r"C:\Users\jnamm\OneDrive\Desktop\Service Cloud for Utilities\Enrollment Center\Enrollment Center - SAP Service Cloud V2 Solution Design v2.0.docx"
doc = Document(SDD)
target = None
for para in doc.paragraphs:
    if para.text.startswith("Before the Enrollment Center is displayed, the eligibility check runs"):
        target = para
        break
assert target is not None, "anchor paragraph not found"
new_p = OxmlElement("w:p")
target._p.addnext(new_p)
np = Paragraph(new_p, target._parent)
run = np.add_run(
    "Customer context header: alongside identification, contract accounts and premise, the header shows the "
    "installed device per utility service (device number, meter number, AMI/AMR type and connection status — "
    "Connected / Comm fault / Manual read / Disconnected) and the customer's DER equipment (EV charger, "
    "battery storage, standby generation) sourced from the DER platform. Both combo account patterns are "
    "supported: one contract account per utility, or a single contract account spanning divisions (joint "
    "invoicing); eligibility is evaluated per contract account and division."
)
run.font.name = "Arial"; run.font.size = Pt(10.5)
np.paragraph_format.space_after = Pt(6)
doc.save(SDD)
print("SDD note inserted")
