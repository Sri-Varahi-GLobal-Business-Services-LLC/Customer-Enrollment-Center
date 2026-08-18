# -*- coding: utf-8 -*-
"""ONE-SHOT patch (do not re-run): deck demo links -> live org Pages URLs + AI/ML methodology page."""
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

DECK = r"C:\Users\jnamm\OneDrive\Desktop\Service Cloud for Utilities\Enrollment Center\Enrollment Center - Client Presentation v1.0.pptx"
BASE = "sri-varahi-global-business-services-llc.github.io/enrollment-center-demo"
ICE = RGBColor(0xCA, 0xDC, 0xFC)
AMBER = RGBColor(0xFF, 0xB3, 0x00)

prs = Presentation(DECK)

REPL = [
 # slide 4 cards (live-first; keep short with ellipsis pattern for sub-pages)
 ("demo.svgbs.com  (pending DNS; live: sri-varahi-global-business-services-llc.github.io/enrollment-center-demo)",
  BASE + "/", 9.5),
 ("demo.svgbs.com/workcenter.html", "\u2026/enrollment-center-demo/workcenter.html", 11.5),
 ("demo.svgbs.com/movein.html", "\u2026/enrollment-center-demo/movein.html", 11.5),
 # closing slide
 ("Demo: demo.svgbs.com   \u00b7   Full design set available for review",
  "Demo: " + BASE + "  (soon: demo.svgbs.com)   \u00b7   Full design set available for review", 11),
]

count = 0
demo_slide = None
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if "See it working today" in shape.text_frame.text:
            demo_slide = slide
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for old, new, size in REPL:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                        run.font.size = Pt(size)
                        count += 1

# add methodology + vanity-domain line at the bottom of the demo slide
if demo_slide is not None:
    box = demo_slide.shapes.add_textbox(Inches(0.6), Inches(6.62), Inches(12.1), Inches(0.4))
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = "\u2728 AI/ML methodology page: \u2026/enrollment-center-demo/ai-roadmap.html"
    r1.font.size = Pt(11.5); r1.font.bold = True; r1.font.name = "Arial"; r1.font.color.rgb = AMBER
    r2 = p.add_run(); r2.text = "    \u00b7    custom domain demo.svgbs.com activates once DNS lands"
    r2.font.size = Pt(11); r2.font.name = "Arial"; r2.font.color.rgb = ICE

prs.save(DECK)
print("url replacements:", count, "| methodology line added:", demo_slide is not None)
