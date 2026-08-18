# -*- coding: utf-8 -*-
"""ONE-SHOT patch (do not re-run): insert 'How the AI score is calculated' slide after the AI slide (pos 12)."""
from pptx import Presentation
from pptx.util import Inches as PIn, Pt as PPt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

DECK = r"C:\Users\jnamm\OneDrive\Desktop\Service Cloud for Utilities\Enrollment Center\Enrollment Center - Client Presentation v1.0.pptx"
NAVY=RGBColor(0x12,0x23,0x3F); BLUE=RGBColor(0x00,0x70,0xF2); BLUED=RGBColor(0x00,0x57,0xD2)
AMBER=RGBColor(0xFF,0xB3,0x00); WHITE=RGBColor(0xFF,0xFF,0xFF); INK=RGBColor(0x1D,0x2D,0x3E)
MUTED=RGBColor(0x55,0x6B,0x82); LIGHT=RGBColor(0xF5,0xF6,0xF7); CARD_B=RGBColor(0xD9,0xD9,0xD9)
ICE=RGBColor(0xCA,0xDC,0xFC); PURPC=RGBColor(0xF0,0xEA,0xFD); PURP=RGBColor(0x6C,0x32,0xA9)

prs = Presentation(DECK)
if any("How the AI score is calculated" in (sh.text_frame.text if sh.has_text_frame else "") for sl in prs.slides for sh in sl.shapes):
    raise SystemExit("deck already has score slide - skipped")

s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE

def tb(x,y,w,h,text,size=14,color=INK,bold=False,align=PP_ALIGN.LEFT,ls=1.0):
    box=s.shapes.add_textbox(PIn(x),PIn(y),PIn(w),PIn(h)); tf=box.text_frame; tf.word_wrap=True
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,ln in enumerate(text.split("\n")):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        if ls!=1.0: p.line_spacing=ls
        r=p.add_run(); r.text=ln; r.font.size=PPt(size); r.font.bold=bold; r.font.name="Arial"; r.font.color.rgb=color
    return box

def card(x,y,w,h,fill=WHITE,border=CARD_B):
    sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,PIn(x),PIn(y),PIn(w),PIn(h))
    sh.adjustments[0]=0.08; sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.color.rgb=border; sh.line.width=PPt(0.75); sh.shadow.inherit=False
    return sh

def numcircle(x,y,n):
    c=s.shapes.add_shape(MSO_SHAPE.OVAL,PIn(x),PIn(y),PIn(0.42),PIn(0.42))
    c.fill.solid(); c.fill.fore_color.rgb=BLUE; c.line.fill.background(); c.shadow.inherit=False
    tf=c.text_frame; tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
    r=p.add_run(); r.text=str(n); r.font.size=PPt(15); r.font.bold=True; r.font.name="Arial"; r.font.color.rgb=WHITE

tb(0.6,0.38,12.1,0.75,"How the AI score is calculated",size=30,bold=True)
tb(0.6,1.02,12.1,0.4,"The percentage is a calibrated enrollment probability — the likelihood this customer enrolls if the agent offers the program",size=14,color=MUTED)

steps=[("Training data","The system labels its own examples: every offer shown, plus the outcome — enrolled within 30 days or not"),
       ("Features","The context already on screen: meter type and status, DER equipment, payment behavior, program history, usage"),
       ("Model","A classification model on SAP HANA Cloud turns those features into a raw probability"),
       ("Calibration","Adjusted so the number means what it says: of customers scored 87%, about 87% actually enroll"),
       ("Reason shown","The line under the badge comes from the model's top contributing factors — never free text"),
       ("Ranking","Sort order weighs probability × program value — high-value programs can outrank easy wins")]
xs=[0.6,4.65,8.7]; ys=[1.75,3.55]; idx=0
for yy in ys:
    for xx in xs:
        h,b=steps[idx]; idx+=1
        card(xx,yy,3.85,1.62,fill=LIGHT)
        numcircle(xx+0.22,yy+0.2,idx)
        tb(xx+0.78,yy+0.24,2.95,0.35,h,size=14,bold=True)
        tb(xx+0.28,yy+0.66,3.35,0.9,b,size=10.5,color=MUTED,ls=1.1)

card(0.6,5.5,12.1,1.3,fill=NAVY,border=NAVY)
tb(0.9,5.68,7.6,0.95,"Cold start: launches on weighted business rules, swaps to the trained model after 3–6 months of outcomes.\nThe model never decides eligibility — that stays in the BRFplus rules engine.",size=12.5,color=ICE,ls=1.2)
tb(8.8,5.68,3.6,0.95,"Live in the demo:\nclick any ✨ AI badge for this explanation",size=12.5,color=AMBER,bold=True,ls=1.2)

# move new slide (last) to position 12 (index 11), right after the AI slide
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
sldIdLst.remove(ids[-1])
sldIdLst.insert(11, ids[-1])
prs.save(DECK)
print("score-calculation slide inserted at position 12; total slides:", len(prs.slides._sldIdLst))
